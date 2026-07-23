# Litigation/utils/litigation_extractors.py
"""
Text and basic extraction for litigation documents.
Supports PDF, DOCX, PPTX, TXT.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

try:
    import fitz  # PyMuPDF for PDFs
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_from_file(file_path: Path):
    """Extract text and list of image paths (if any) from supported file types."""
    suffix = file_path.suffix.lower()
    text = ""
    images = []

    try:
        if suffix == ".pdf" and fitz:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text("text") + "\n"
                # Extract images from PDF pages
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_path = file_path.parent / f"{file_path.stem}_img_{page.number}_{img_index}.png"
                    image_path.write_bytes(image_bytes)
                    images.append(image_path)
            doc.close()

        elif suffix in (".docx", ".doc") and Document:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        elif suffix in (".pptx", ".ppt") and Presentation:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        elif suffix == ".txt":
            text = file_path.read_text(encoding="utf-8", errors="replace")

        else:
            logger.warning(f"Unsupported file type for extraction: {file_path.suffix}")

    except Exception as e:
        logger.error(f"Extraction failed for {file_path.name}: {e}")

    return text.strip(), images